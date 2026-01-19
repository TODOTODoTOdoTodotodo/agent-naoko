import os
from pptx import Presentation
from rich.console import Console
from pathlib import Path

console = Console()

class DocParser:
    @staticmethod
    def parse(file_path: str) -> str:
        """
        Parses the document based on its extension and returns raw text.
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.pptx':
            return DocParser._parse_pptx(file_path)
        elif ext == '.pdf':
            # Placeholder for PDF
            return f"[Parser] PDF parsing not yet implemented for {file_path}"
        elif ext in ['.xlsx', '.xls']:
            # Placeholder for Excel
            return f"[Parser] Excel parsing not yet implemented for {file_path}"
        elif ext == '.md':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            return f"[Parser] Unsupported file format: {ext}"

    @staticmethod
    def _parse_pptx(file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            full_text = []
            
            console.print(f"[dim][DocParser] Extracting text from {Path(file_path).name}...[/dim]")

            for i, slide in enumerate(prs.slides):
                slide_text = []
                # Extract title
                if slide.shapes.title:
                    slide_text.append(f"## Slide {i+1}: {slide.shapes.title.text}")
                else:
                    slide_text.append(f"## Slide {i+1}")

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Avoid duplicating title
                        if shape == slide.shapes.title:
                            continue
                        clean_text = shape.text.replace('\n', ' ').strip()
                        if clean_text:
                            slide_text.append(f"- {clean_text}")
                
                full_text.append("\n".join(slide_text))
            
            return "\n\n".join(full_text)
            
        except Exception as e:
            console.print(f"[red][DocParser] Error parsing PPTX: {e}[/red]")
            return ""
